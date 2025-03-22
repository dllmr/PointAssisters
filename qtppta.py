# /// script
# requires-python = ">=3.12"
# dependencies = [
#     "markdown",
#     "matplotlib",
#     "python-pptx",
#     "PySide6",
# ]
# ///


import logging
import markdown
import sys
from collections import defaultdict
from matplotlib import font_manager as fm
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.text.text import _Paragraph
from PySide6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, 
                            QHBoxLayout, QPushButton, QFileDialog, QLineEdit, 
                            QTextEdit, QCheckBox, QGroupBox, QStatusBar, QLabel)
from typing import Dict, List, Set, Tuple, Any
from xml.etree import ElementTree
from pptx.opc.constants import RELATIONSHIP_TYPE as RT


INTERNAL_FONT_MARKERS = frozenset({
    '+mj-lt',    # Default Latin font for gothic text
    '+mn-lt',    # Default Latin font for mincho text
    '+body',     # Default body font
    '+major',    # Default major font
    '+minor',    # Default minor font
    '@',         # Font fallback marker
})

THEME_FONT_CODES = {
    '+mj-lt': 'Major Latin',
    '+mn-lt': 'Minor Latin',
    '+mj-ea': 'Major East Asian',
    '+mn-ea': 'Minor East Asian',
    '+mj-cs': 'Major Complex Script',
    '+mn-cs': 'Minor Complex Script',
    '+mj-sym': 'Major Symbol',
    '+mn-sym': 'Minor Symbol',
}

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(levelname)s:%(name)s:%(message)s'
)

logger = logging.getLogger(__name__)

def find_hidden_slides(pptx_path: str) -> List[int]:
    prs = Presentation(pptx_path)
    hidden_slides = []
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        try:
            # Check if slide is marked as hidden
            if hasattr(slide, '_element') and slide._element.get('show') == '0':
                hidden_slides.append(slide_num)
        except Exception as e:
            logger.warning(f"Error checking slide {slide_num}: {e}")
            
    return hidden_slides

def count_words_in_shape(shape: BaseShape) -> int:
    """Count the words in a PowerPoint shape."""
    word_count = 0
    
    try:
        # Handle text frames
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            # Count words in text frame paragraphs
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    word_count += len(paragraph.text.split())
                
        # Handle tables
        if hasattr(shape, 'has_table') and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        if paragraph.text.strip():
                            word_count += len(paragraph.text.split())
                    
    except Exception as e:
        logger.debug(f"Error counting words in shape: {str(e)}")
        
    return word_count

def analyze_presentation_statistics(pptx_path: str) -> Dict[str, Any]:
    """Analyze general statistics about the presentation."""
    prs = Presentation(pptx_path)
    stats = {
        "total_slides": len(prs.slides),
        "hidden_slides": [],
        "total_words": 0,
        "slide_word_counts": {},
        "max_words_slide": 0,
        "max_words_count": 0
    }
    
    # Find hidden slides and count words per slide
    for slide_num, slide in enumerate(prs.slides, start=1):
        try:
            # Check if slide is hidden
            if hasattr(slide, '_element') and slide._element.get('show') == '0':
                stats["hidden_slides"].append(slide_num)
            
            # Count words on this slide
            slide_word_count = 0
            for shape in slide.shapes:
                slide_word_count += count_words_in_shape(shape)
            
            stats["slide_word_counts"][slide_num] = slide_word_count
            stats["total_words"] += slide_word_count
            
            # Track slide with most words
            if slide_word_count > stats["max_words_count"]:
                stats["max_words_count"] = slide_word_count
                stats["max_words_slide"] = slide_num
            
        except Exception as e:
            logger.warning(f"Error analyzing slide {slide_num}: {e}")
    
    return stats

def generate_presentation_summary(pptx_path: str) -> str:
    """Generate a summary section with general presentation statistics."""
    stats = analyze_presentation_statistics(pptx_path)
    
    # Get just the filename without the full path
    filename = Path(pptx_path).name
    
    result = f"## Presentation Summary for {filename}\n"
    
    # Basic stats
    result += f"Total slides: {stats['total_slides']}<br />\n"
    
    # Hidden slides count
    hidden_count = len(stats["hidden_slides"])
    if hidden_count > 0:
        result += f"Hidden slides: {hidden_count} ({', '.join(str(num) for num in sorted(stats['hidden_slides']))})<br />\n"
    else:
        result += "Hidden slides: 0<br />\n"
    
    # Word counts
    result += f"Total words: {stats['total_words']}<br />\n"
    
    if stats["max_words_count"] > 0:
        result += f"Slide with most words: {stats['max_words_slide']} ({stats['max_words_count']} words)<br />\n"
    
    result += "***\n"
    return result

def generate_hidden_slides_report(pptx_path: str) -> str:
    hidden_slides = find_hidden_slides(pptx_path)
    
    result = ""

    result += "## Hidden Slides\n"
    
    if hidden_slides:
        result += "Hidden slides: " + (", ".join(str(num) for num in sorted(hidden_slides))) + "\n"
    else:
        result += "(no hidden slides found)\n"
    result += "***\n"

    return result

def find_animations_and_transitions(pptx_path: str) -> Tuple[Set[int], Set[int]]:
    """
    Find slides containing transitions or animations in a PowerPoint presentation.
    
    Args:
        pptx_path: Path to the PowerPoint file
        
    Returns:
        Tuple containing:
        - Set of slide numbers with transitions
        - Set of slide numbers with animations
    """
    prs = Presentation(pptx_path)
    slides_with_transitions = set()
    slides_with_animations = set()
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        try:
            # Check for transitions
            transition = slide._element.find('./p:transition', 
                                          {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
            if transition is not None:
                slides_with_transitions.add(slide_num)
            
            # Check for animations
            timing = slide._element.find('./p:timing',
                                      {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
            if timing is not None:
                # Look for any animation elements
                anim_elements = timing.findall('.//p:anim',
                                            {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
                anim_elements.extend(timing.findall('.//p:animEffect',
                                                  {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}))
                if anim_elements:
                    slides_with_animations.add(slide_num)
                    
        except Exception as e:
            logger.warning(f"Error processing slide {slide_num}: {e}")
            continue
            
    return slides_with_transitions, slides_with_animations

def format_effects_report(slides_with_transitions: Set[int], slides_with_animations: Set[int]) -> str:
    result = ""

    result += "## Transitions and Animations\n"
    
    if slides_with_transitions:
        result += "Slides with transitions: " + (", ".join(str(num) for num in sorted(slides_with_transitions))) + "<br />\n"
    else:
        result += "(no transitions found)<br />\n"
        
    if slides_with_animations:
        result += "Slides with animations: " + (", ".join(str(num) for num in sorted(slides_with_animations))) + "\n"
    else:
        result += "(no animations found)\n"
    result += "***\n"

    return result

def generate_effects_report(pptx_path: str) -> str:
    transitions, animations = find_animations_and_transitions(pptx_path)
    return format_effects_report(transitions, animations)
    
def get_system_fonts() -> Set[str]:
    """
    Get a set of all system fonts including both TTF and OTF formats.
    
    Returns:
        A sorted set of font names available on the system.
    """
    font_names: List[str] = []
    
    fonts = fm.findSystemFonts(fontpaths=None)
    
    # Process each font file to get its name
    for font in fonts:
        try:
            # Attempt to get the font properties
            font_name = fm.FontProperties(fname=font).get_name()
            font_names.append(font_name)
        except Exception as e:
            # Log debug info about font loading errors
            logger.debug(f"Error loading font properties for {font}: {e}")

    # Return sorted unique font names
    return sorted(set(font_names))

def analyze_paragraph_fonts(paragraph: _Paragraph) -> Dict[str, Dict[str, Any]]:
    """
    Extract fonts from a paragraph, including runs.
    
    Returns:
        Dict mapping font names to info containing:
        - has_visible_text: Whether the font contains visible text (True) or only whitespace (False)
        - sizes: Set of font sizes used with visible text
    """
    fonts = {}
    
    # Track runs with missing font information but with size/text data
    unknown_sizes = set()
    has_unknown_visible_text = False

    for run in paragraph.runs:
        try:
            # Check if this run contains non-whitespace characters
            has_visible_text = bool(run.text.strip())
            
            # Get font size if available
            font_size = None
            if hasattr(run.font, 'size') and run.font.size is not None:
                # Convert from EMUs to points (1 point = 12700 EMUs)
                if isinstance(run.font.size, int):
                    font_size = int(round(run.font.size / 12700))
            
            if hasattr(run, 'font') and run.font.name and not is_internal_font(run.font.name):
                # Process normal font with name
                font_name = run.font.name
                
                # Initialize font info if not already in dictionary
                if font_name not in fonts:
                    fonts[font_name] = {
                        "has_visible_text": False,
                        "sizes": set()
                    }
                
                # Update visibility status
                if has_visible_text:
                    fonts[font_name]["has_visible_text"] = True
                    
                    # Only track size if it's available and this run has visible text
                    if font_size is not None:
                        fonts[font_name]["sizes"].add(font_size)
            elif has_visible_text or font_size:
                # Track that we have text/size but no font info
                has_unknown_visible_text = has_unknown_visible_text or has_visible_text
                if font_size is not None and has_visible_text:
                    unknown_sizes.add(font_size)

        except Exception as e:
            logger.debug(f"Error analyzing run: {str(e)}")
    
    # Add unknown font entry if we found runs with missing font information
    if has_unknown_visible_text or unknown_sizes:
        fonts["(unknown)"] = {
            "has_visible_text": has_unknown_visible_text,
            "sizes": unknown_sizes
        }

    return fonts

def analyze_shape_fonts(shape: BaseShape) -> Dict[str, Dict[str, Any]]:
    """
    Safely extract fonts from a shape, tracking whether each font has visible text and its sizes.
    
    Returns:
        Dict mapping font names to info containing:
        - has_visible_text: Whether the font contains visible text (True) or only whitespace (False)
        - sizes: Set of font sizes used with visible text
    """
    fonts = {}
    
    try:
        # Handle text frames
        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                paragraph_fonts = analyze_paragraph_fonts(paragraph)
                # Merge results, keeping track of visible text status and sizes
                for font_name, font_info in paragraph_fonts.items():
                    if font_name not in fonts:
                        fonts[font_name] = {
                            "has_visible_text": False,
                            "sizes": set()
                        }
                    
                    # Update visibility
                    fonts[font_name]["has_visible_text"] = fonts[font_name]["has_visible_text"] or font_info["has_visible_text"]
                    
                    # Add sizes if this font has visible text
                    if font_info["has_visible_text"]:
                        fonts[font_name]["sizes"].update(font_info["sizes"])
                
        # Handle tables
        if hasattr(shape, 'has_table') and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph_fonts = analyze_paragraph_fonts(paragraph)
                        # Merge results, keeping track of visible text status and sizes
                        for font_name, font_info in paragraph_fonts.items():
                            if font_name not in fonts:
                                fonts[font_name] = {
                                    "has_visible_text": False,
                                    "sizes": set()
                                }
                            
                            # Update visibility
                            fonts[font_name]["has_visible_text"] = fonts[font_name]["has_visible_text"] or font_info["has_visible_text"]
                            
                            # Add sizes if this font has visible text
                            if font_info["has_visible_text"]:
                                fonts[font_name]["sizes"].update(font_info["sizes"])
        
        # Handle group shapes - recursively process shapes within groups
        if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            if hasattr(shape, 'shapes'):
                for child_shape in shape.shapes:
                    child_fonts = analyze_shape_fonts(child_shape)
                    # Merge results
                    for font_name, font_info in child_fonts.items():
                        if font_name not in fonts:
                            fonts[font_name] = {
                                "has_visible_text": False,
                                "sizes": set()
                            }
                        
                        # Update visibility
                        fonts[font_name]["has_visible_text"] = fonts[font_name]["has_visible_text"] or font_info["has_visible_text"]
                        
                        # Add sizes if this font has visible text
                        if font_info["has_visible_text"]:
                            fonts[font_name]["sizes"].update(font_info["sizes"])
                        
    except Exception as e:
        logger.debug(f"Error analyzing shape: {str(e)}")
        
    return fonts

def analyze_fonts(pptx_path: str) -> Tuple[Dict[int, Dict[str, Dict[str, Any]]], Dict[str, Dict[str, Any]]]:
    """
    Analyze fonts used in a PowerPoint presentation.
    
    Returns:
        Tuple containing:
        - Dictionary mapping slide numbers to shape types to font usage info
        - Dictionary mapping all fonts to their visibility and size information
    """
    prs = Presentation(pptx_path)
    font_usage = defaultdict(lambda: defaultdict(dict))
    all_fonts_info = {}
    
    for slide_num, slide in enumerate(prs.slides, start=1):
        try:
            for shape in slide.shapes:
                try:
                    shape_type = f"Text Shape: {shape.name}" if hasattr(shape, 'name') else "Shape"
                    fonts = analyze_shape_fonts(shape)
                    
                    if fonts:
                        font_usage[slide_num][shape_type] = fonts
                        
                        # Update global font tracking
                        for font_name, font_info in fonts.items():
                            if font_name not in all_fonts_info:
                                all_fonts_info[font_name] = {
                                    "has_visible_text": False,
                                    "sizes": set()
                                }
                            
                            # Update visibility
                            all_fonts_info[font_name]["has_visible_text"] = (
                                all_fonts_info[font_name]["has_visible_text"] or font_info["has_visible_text"]
                            )
                            
                            # Add sizes if this font has visible text
                            if font_info["has_visible_text"]:
                                all_fonts_info[font_name]["sizes"].update(font_info["sizes"])
                        
                except Exception as e:
                    logger.debug(f"Error processing shape in slide {slide_num}: {str(e)}")
                    continue
                    
        except Exception as e:
            logger.warning(f"Error processing slide {slide_num}: {str(e)}")
            continue

    return font_usage, all_fonts_info

def is_internal_font(font_name: str) -> bool:
    if not font_name:
        return True
    font_name = font_name.lower()
    return any(font_name.startswith(marker) for marker in INTERNAL_FONT_MARKERS)

def extract_theme_fonts(presentation: Any) -> Dict[str, Any]:
    """Extract theme font information from the presentation."""
    theme_fonts = {}
    
    try:
        # Get the default theme from the first slide master
        if presentation.slide_masters and len(presentation.slide_masters) > 0:
            master = presentation.slide_masters[0]
            
            # Access the theme through the part relationships
            try:
                # Get the master part
                master_part = master.part
                
                # Find theme relationships
                theme_rels = [rel for rel in master_part.rels.values() 
                             if rel.reltype == RT.THEME]
                
                if theme_rels:
                    # Get the first theme part using the relationship ID
                    theme_rel = theme_rels[0]
                    theme_part = master_part.related_part(theme_rel.rId)
                    
                    # Parse the theme XML
                    theme_element = ElementTree.fromstring(theme_part.blob)
                    
                    # Extract font scheme
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    font_scheme_elem = theme_element.find('.//a:fontScheme', ns)
                    
                    if font_scheme_elem is not None:
                        # Get font scheme name
                        scheme_name = font_scheme_elem.get('name', 'Unknown')
                        
                        # Get major font element
                        major_font_elem = font_scheme_elem.find('.//a:majorFont', ns)
                        major_fonts = {}
                        
                        if major_font_elem is not None:
                            latin = major_font_elem.find('.//a:latin', ns)
                            ea = major_font_elem.find('.//a:ea', ns)
                            cs = major_font_elem.find('.//a:cs', ns)
                            sym = major_font_elem.find('.//a:sym', ns)
                            
                            major_fonts = {
                                "latin": latin.get('typeface') if latin is not None else None,
                                "east_asian": ea.get('typeface') if ea is not None else None,
                                "complex_script": cs.get('typeface') if cs is not None else None,
                                "symbol": sym.get('typeface') if sym is not None else None
                            }
                        
                        # Get minor font element
                        minor_font_elem = font_scheme_elem.find('.//a:minorFont', ns)
                        minor_fonts = {}
                        
                        if minor_font_elem is not None:
                            latin = minor_font_elem.find('.//a:latin', ns)
                            ea = minor_font_elem.find('.//a:ea', ns)
                            cs = minor_font_elem.find('.//a:cs', ns)
                            sym = minor_font_elem.find('.//a:sym', ns)
                            
                            minor_fonts = {
                                "latin": latin.get('typeface') if latin is not None else None,
                                "east_asian": ea.get('typeface') if ea is not None else None,
                                "complex_script": cs.get('typeface') if cs is not None else None,
                                "symbol": sym.get('typeface') if sym is not None else None
                            }
                        
                        theme_fonts = {
                            "scheme_name": scheme_name,
                            "major_fonts": major_fonts,
                            "minor_fonts": minor_fonts
                        }
            except Exception as e:
                theme_fonts["error"] = str(e)
    except Exception as e:
        theme_fonts["error"] = str(e)
    
    return theme_fonts

def format_font_report(font_usage: Dict[int, Dict[str, Dict[str, Any]]], 
                     all_fonts_info: Dict[str, Dict[str, Any]],
                     system_fonts: Set[str],
                     presentation: Any,
                     font_size_threshold: int = 24) -> str:
    """Create a formatted report showing font usage and theme fonts."""
    result = ""

    # Add CSS styling for tables
    result += """<style>
table {
    border-collapse: collapse;
    margin: 10px 0;
    width: 100%;
}
th, td {
    padding: 8px 16px;
    text-align: left;
    border: 1px solid gray;
}
th {
    font-weight: bold;
}
.whitespace-only {
    font-style: italic;
}
.small-font {
    font-weight: bold;
}
.unknown-font {
    font-style: italic;
}
.unknown-small-font {
    font-weight: bold;
    font-style: italic;
}
</style>
"""

    # Filter out internal fonts
    all_fonts_info = {f.strip(): v for f, v in all_fonts_info.items() if not is_internal_font(f)}
    
    # Create a normalized version of system fonts for flexible matching
    system_fonts_lower = {s.lower().strip() for s in system_fonts}
    
    # Create a mapping with normalized versions (no spaces, no punctuation)
    normalized_system_fonts = {}
    for font in system_fonts:
        # Create normalized version (lowercase, no spaces, no punctuation)
        normalized = ''.join(c.lower() for c in font if c.isalnum())
        normalized_system_fonts[normalized] = font
    
    # Create a mapping of fonts to the slides that use them, with visibility information
    font_to_slides: Dict[str, Dict[int, Dict[str, Any]]] = {}
    for slide_num, shapes in font_usage.items():
        # Process all fonts from all shapes in this slide
        for shape_type, fonts_info in shapes.items():
            for font, font_info in fonts_info.items():
                if font and not is_internal_font(font):
                    if font not in font_to_slides:
                        font_to_slides[font] = {}
                    
                    # Store font info for this slide
                    if slide_num not in font_to_slides[font]:
                        font_to_slides[font][slide_num] = {
                            "has_visible_text": False,
                            "sizes": set()
                        }
                    
                    # Update visibility for this slide
                    font_to_slides[font][slide_num]["has_visible_text"] = (
                        font_to_slides[font][slide_num]["has_visible_text"] or font_info["has_visible_text"]
                    )
                    
                    # Add sizes if this font has visible text on this slide
                    if font_info["has_visible_text"]:
                        font_to_slides[font][slide_num]["sizes"].update(font_info["sizes"])

    result += "## Custom Font Usage\n"
    
    if font_to_slides:
        result += "<table>\n"
        result += "<tr><th>Font Name</th><th>Local Status</th><th>Used on Slides</th><th>Font Sizes</th><th>Notes</th></tr>\n"
        
        # Process fonts by moving unknown to the end but otherwise alphabetical
        sorted_fonts = sorted(
            font_to_slides.keys(), 
            key=lambda x: (x == "(unknown)", x.lower())
        )
        
        for font in sorted_fonts:
            if not font:  # Skip None values
                continue
                
            # Special handling for unknown fonts
            is_unknown = font == "(unknown)"
            font_name_display = f"<span class='unknown-font'>{font}</span>" if is_unknown else font
            
            # Determine font status with flexible matching
            if is_unknown:
                status = "<span class='unknown-font'>Unknown (theme/default font)</span>"
            else:
                # First try exact match
                if font.lower() in system_fonts_lower:
                    status = "✅ Installed"
                else:
                    # Try normalized matching (PowerPoint-like flexibility)
                    font_normalized = ''.join(c.lower() for c in font if c.isalnum())
                    if font_normalized in normalized_system_fonts:
                        matched_font = normalized_system_fonts[font_normalized]
                        status = f"✅ Installed (as '{matched_font}')"
                    else:
                        status = "❌ Missing"
                
            # Convert slide numbers to a readable string, marking whitespace-only slides
            slides_info = sorted(font_to_slides[font].items())
            slide_parts = []
            
            # Track all sizes for this font
            all_sizes = set()
            
            # Track slides with small fonts (below threshold)
            small_font_slides = set()
            
            for slide_num, info in slides_info:
                if info["has_visible_text"]:
                    # Check for small fonts on this slide
                    has_small_font = any(size < font_size_threshold for size in info["sizes"])
                    
                    if has_small_font:
                        small_font_slides.add(slide_num)
                        slide_parts.append(f"<span class='small-font'>{slide_num}†</span>")
                    else:
                        slide_parts.append(str(slide_num))
                        
                    # Collect sizes from visible text
                    all_sizes.update(info["sizes"])
                else:
                    slide_parts.append(f"<span class='whitespace-only'>{slide_num}*</span>")
            
            slides_str = ", ".join(slide_parts)
            
            # Format sizes as a sorted list
            sizes_str = ""
            if all_sizes:
                # Sort sizes and highlight those below threshold
                size_parts = []
                for size in sorted(all_sizes):
                    if size < font_size_threshold:
                        # Use combined style for unknown small fonts
                        if is_unknown:
                            size_parts.append(f"<span class='unknown-small-font'>{size}</span>")
                        else:
                            size_parts.append(f"<span class='small-font'>{size}</span>")
                    else:
                        if is_unknown:
                            size_parts.append(f"<span class='unknown-font'>{size}</span>")
                        else:
                            size_parts.append(str(size))
                sizes_str = ", ".join(size_parts)
            
            # Determine if this font is only used for whitespace across all slides
            whitespace_only = not any(info["has_visible_text"] for _, info in slides_info)
            
            # Add notes about whitespace usage, small fonts, and unknown fonts
            notes = []
            if is_unknown:
                notes.append("<span class='unknown-font'>Font information not available (likely theme or default font)</span>")
            
            if whitespace_only:
                notes.append("<span class='whitespace-only'>Used only for whitespace</span>")
            elif any(not info["has_visible_text"] for _, info in slides_info):
                notes.append("<span class='whitespace-only'>* = whitespace only on marked slides</span>")
            
            if small_font_slides:
                slides_list = ", ".join(str(slide) for slide in sorted(small_font_slides))
                notes.append(f"<span class='small-font'>† = Small font (&lt;{font_size_threshold}pt) on slides {slides_list}</span>")
            
            result += f"<tr><td>{font_name_display}</td><td>{status}</td><td>{slides_str}</td><td>{sizes_str}</td><td>{' '.join(notes)}</td></tr>\n"
        
        result += "</table>\n"
    else:
        result += "(no custom fonts used in presentation)\n"
    
    # Print theme fonts
    result += "\n## Theme Fonts\n"
    
    # Extract theme fonts
    theme_fonts = extract_theme_fonts(presentation)
    
    if "error" in theme_fonts:
        result += f"Error accessing theme fonts: {theme_fonts['error']}\n"
    else:
        if theme_fonts.get("scheme_name"):
            result += f"Theme scheme name: {theme_fonts['scheme_name']}\n\n"
        
        result += "<table>\n"
        result += "<tr><th>Font Type</th><th>Font Name</th><th>Local Status</th></tr>\n"
        
        # Process major fonts
        major_fonts = theme_fonts.get("major_fonts", {})
        for script, font in major_fonts.items():
            if font:
                # Use flexible font matching for theme fonts too
                if font.lower() in system_fonts_lower:
                    status = "✅ Installed"
                else:
                    # Try normalized matching
                    font_normalized = ''.join(c.lower() for c in font if c.isalnum())
                    if font_normalized in normalized_system_fonts:
                        matched_font = normalized_system_fonts[font_normalized]
                        status = f"✅ Installed (as '{matched_font}')"
                    else:
                        status = "❌ Missing"
                result += f"<tr><td>Major {script.replace('_', ' ').title()}</td><td>{font}</td><td>{status}</td></tr>\n"
        
        # Process minor fonts
        minor_fonts = theme_fonts.get("minor_fonts", {})
        for script, font in minor_fonts.items():
            if font:
                # Use flexible font matching for theme fonts too
                if font.lower() in system_fonts_lower:
                    status = "✅ Installed"
                else:
                    # Try normalized matching
                    font_normalized = ''.join(c.lower() for c in font if c.isalnum())
                    if font_normalized in normalized_system_fonts:
                        matched_font = normalized_system_fonts[font_normalized]
                        status = f"✅ Installed (as '{matched_font}')"
                    else:
                        status = "❌ Missing"
                result += f"<tr><td>Minor {script.replace('_', ' ').title()}</td><td>{font}</td><td>{status}</td></tr>\n"
        
        result += "</table>\n"
        
        if not (major_fonts or minor_fonts):
            result += "(no theme fonts defined)\n"
    
    # Print summary statistics
    total_fonts = len([font for font in font_to_slides.keys() if font != "(unknown)"])
    
    # Update missing fonts count to use normalized matching
    missing_fonts = 0
    for font in font_to_slides:
        if font != "(unknown)":
            font_lower = font.lower()
            font_normalized = ''.join(c.lower() for c in font if c.isalnum())
            if font_lower not in system_fonts_lower and font_normalized not in normalized_system_fonts:
                missing_fonts += 1
    
    unknown_fonts = 1 if "(unknown)" in font_to_slides else 0
    
    # Count whitespace-only fonts (excluding unknown)
    whitespace_only_fonts = sum(
        1 for font in font_to_slides 
        if font != "(unknown)" and not any(info["has_visible_text"] for _, info in font_to_slides[font].items())
    )
    
    # Count fonts with sizes below threshold (excluding unknown)
    small_fonts = sum(
        1 for font in font_to_slides
        if font != "(unknown)" and any(
            any(size < font_size_threshold for size in info["sizes"])
            for _, info in font_to_slides[font].items()
            if info["has_visible_text"]  # Only consider visible text
        )
    )
    
    # Count slides with small fonts (including those from unknown fonts)
    slides_with_small_fonts = set()
    for font, slides_info in font_to_slides.items():
        for slide_num, info in slides_info.items():
            if info["has_visible_text"] and any(size < font_size_threshold for size in info["sizes"]):
                slides_with_small_fonts.add(slide_num)
    
    # Count slides with unknown fonts
    slides_with_unknown_fonts = set()
    if "(unknown)" in font_to_slides:
        for slide_num, info in font_to_slides["(unknown)"].items():
            if info["has_visible_text"]:
                slides_with_unknown_fonts.add(slide_num)
    
    # Update theme fonts count with flexible matching
    total_theme_fonts = sum(
        len([f for f in fonts.values() if f]) 
        for fonts in [theme_fonts.get("major_fonts", {}), theme_fonts.get("minor_fonts", {})]
    )
    
    # Update missing theme fonts count with flexible matching
    missing_theme_fonts = 0
    for fonts_dict in [theme_fonts.get("major_fonts", {}), theme_fonts.get("minor_fonts", {})]:
        for font in fonts_dict.values():
            if font:
                font_lower = font.lower()
                font_normalized = ''.join(c.lower() for c in font if c.isalnum())
                if font_lower not in system_fonts_lower and font_normalized not in normalized_system_fonts:
                    missing_theme_fonts += 1
    
    result += "\n## Fonts Summary\n"
    result += f"Total custom fonts: {total_fonts}<br />\n"
    result += f"Missing custom fonts: {missing_fonts}<br />\n"
    
    # Only show whitespace-only fonts line if there are any
    if whitespace_only_fonts > 0:
        result += f"Fonts used only for whitespace: {whitespace_only_fonts}<br />\n"
    
    if unknown_fonts > 0:
        slide_list = ", ".join(str(num) for num in sorted(slides_with_unknown_fonts))
        result += f"<span class='unknown-font'>Unknown fonts (theme/default): {unknown_fonts} (on slides {slide_list})</span><br />\n"
    
    if small_fonts > 0:
        slide_list = ", ".join(str(num) for num in sorted(slides_with_small_fonts))
        result += f"<span class='small-font'>Fonts below {font_size_threshold}pt: {small_fonts} (on slides {slide_list})</span><br />\n"
    
    # Add note about small font sizes in unknown fonts
    unknown_small_fonts = False
    if "(unknown)" in font_to_slides:
        for _, info in font_to_slides["(unknown)"].items():
            if info["has_visible_text"] and any(size < font_size_threshold for size in info["sizes"]):
                unknown_small_fonts = True
                break
    
    if unknown_small_fonts:
        result += f"<span class='unknown-small-font'>Note: Small font sizes detected in unknown fonts</span><br />\n"
    
    result += f"Total theme fonts: {total_theme_fonts}<br />\n"
    result += f"Missing theme fonts: {missing_theme_fonts}\n"
    result += "***\n"

    return result

def generate_font_report(pptx_path: str, font_size_threshold: int) -> str:
    # Get system fonts
    system_fonts = get_system_fonts()
    
    # Open presentation
    prs = Presentation(pptx_path)
    
    # Analyze presentation
    font_usage, all_fonts_info = analyze_fonts(pptx_path)
    
    # Format report
    return format_font_report(font_usage, all_fonts_info, system_fonts, prs, font_size_threshold)

class PowerPointAnalyzerGUI(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("PointAssisters - PowerPoint Analyzer")
        self.resize(800, 600)

        # Create central widget and main layout
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        layout = QVBoxLayout(central_widget)

        # File selection
        file_layout = QHBoxLayout()
        self.file_entry = QLineEdit()
        self.file_entry.setPlaceholderText("Select a PowerPoint file...")
        browse_button = QPushButton("Browse")
        browse_button.clicked.connect(self.browse_file)
        file_layout.addWidget(self.file_entry)
        file_layout.addWidget(browse_button)
        layout.addLayout(file_layout)

        # Analysis options
        options_group = QGroupBox("Analysis Options")
        options_layout = QHBoxLayout()
        self.summary_check = QCheckBox("Summary")
        self.hidden_check = QCheckBox("Hidden Slides")
        self.effects_check = QCheckBox("Effects")
        self.fonts_check = QCheckBox("Fonts")
        self.summary_check.setChecked(True)
        self.hidden_check.setChecked(True)
        self.effects_check.setChecked(True)
        self.fonts_check.setChecked(True)
        options_layout.addWidget(self.summary_check)
        options_layout.addWidget(self.hidden_check)
        options_layout.addWidget(self.effects_check)
        options_layout.addWidget(self.fonts_check)
        options_group.setLayout(options_layout)
        layout.addWidget(options_group)
        
        # Font size threshold
        threshold_layout = QHBoxLayout()
        threshold_layout.addWidget(QWidget())  # Spacer
        threshold_layout.addWidget(QLabel("Font Size Threshold:"))
        self.threshold_entry = QLineEdit()
        self.threshold_entry.setPlaceholderText("24")
        self.threshold_entry.setMaximumWidth(50)
        self.threshold_entry.setToolTip("Font sizes below this value will be flagged (default: 24)")
        threshold_layout.addWidget(self.threshold_entry)
        threshold_layout.addWidget(QLabel("points"))
        threshold_layout.addStretch(1)  # Add stretch to push widgets to the left
        layout.addLayout(threshold_layout)

        # Analyze button
        analyze_button = QPushButton("Analyze")
        analyze_button.clicked.connect(self.analyze)
        layout.addWidget(analyze_button)

        # Results area
        self.results_text = QTextEdit()
        self.results_text.setReadOnly(True)
        layout.addWidget(self.results_text)

        # Status bar
        self.status_bar = QStatusBar()
        self.setStatusBar(self.status_bar)

    def browse_file(self):
        filename, _ = QFileDialog.getOpenFileName(
            self,
            "Select PowerPoint File",
            "",
            "PowerPoint files (*.pptx);;All files (*.*)"
        )
        if filename:
            self.file_entry.setText(filename)
            self.results_text.clear()
            self.status_bar.clearMessage()

    def analyze(self):
        file_path = self.file_entry.text()
        if not file_path:
            self.status_bar.showMessage("Please select a file first")
            return

        if not Path(file_path).exists():
            self.status_bar.showMessage("Selected file does not exist")
            return

        # Get font size threshold (default to 24 if not provided or invalid)
        try:
            font_size_threshold = int(self.threshold_entry.text())
            if font_size_threshold <= 0:
                font_size_threshold = 24
        except (ValueError, TypeError):
            font_size_threshold = 24

        self.results_text.clear()
        self.status_bar.showMessage("Analyzing...")
        QApplication.processEvents()  # Ensure UI updates

        try:
            # Capture output
            output = ""
            
            # Include the presentation summary first if selected
            if self.summary_check.isChecked():
                output += generate_presentation_summary(file_path)
            
            # Add other selected analysis sections
            if self.hidden_check.isChecked():
                output += generate_hidden_slides_report(file_path)
            if self.effects_check.isChecked():
                output += generate_effects_report(file_path)
            if self.fonts_check.isChecked():
                output += generate_font_report(file_path, font_size_threshold)

            # Display results
            html = markdown.markdown(output)
            self.results_text.setHtml(html)
            self.status_bar.showMessage("Analysis complete")

        except Exception as e:
            self.status_bar.showMessage(f"Error: {str(e)}")
            logger.error(f"Error analyzing presentation: {e}")

def main():
    app = QApplication(sys.argv)
    window = PowerPointAnalyzerGUI()
    window.show()
    sys.exit(app.exec())

if __name__ == "__main__":
    main()
