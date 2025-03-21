# /// script
# requires-python = ">=3.12"
# dependencies = [
#     "python-pptx",
# ]
# ///


import json
import sys
from pathlib import Path
from pptx import Presentation
from typing import Any, Dict, List
from xml.etree import ElementTree


def resolve_theme_font(shape: Any, theme_code: str) -> str:
    """Resolve theme font codes to actual font names."""
    try:
        if not theme_code:
            return None
            
        # If it's not a theme code (doesn't start with +), return as is
        if not theme_code.startswith('+'):
            return theme_code
            
        # Get access to the theme
        theme_part = None
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        
        if hasattr(shape, 'part') and hasattr(shape.part, 'slide'):
            # For shapes on slides
            slide_layout = shape.part.slide.slide_layout
            master = slide_layout.slide_master
            master_part = master.part  # part is on the SlideMaster object, not on _element
            theme_rels = [rel for rel in master_part.rels.values() if rel.reltype == RT.THEME]
            if theme_rels:
                theme_rel = theme_rels[0]
                theme_part = master_part.related_part(theme_rel.rId)
        elif hasattr(shape, 'slide_layout') and hasattr(shape.slide_layout, 'slide_master'):
            # For slides
            master = shape.slide_layout.slide_master
            master_part = master.part  # part is on the SlideMaster object, not on _element
            theme_rels = [rel for rel in master_part.rels.values() if rel.reltype == RT.THEME]
            if theme_rels:
                theme_rel = theme_rels[0]
                theme_part = master_part.related_part(theme_rel.rId)
        elif hasattr(shape, '_element') and hasattr(shape, 'part'):
            # For slide masters
            master_part = shape.part  # part is on the SlideMaster object, not on _element
            theme_rels = [rel for rel in master_part.rels.values() if rel.reltype == RT.THEME]
            if theme_rels:
                theme_rel = theme_rels[0]
                theme_part = master_part.related_part(theme_rel.rId)
        
        if not theme_part:
            return f"Unable to resolve theme code: {theme_code} (no theme found)"
            
        # Parse the theme XML
        from xml.etree import ElementTree
        theme_element = ElementTree.fromstring(theme_part.blob)
        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        
        # Find font scheme
        font_scheme = theme_element.find('.//a:fontScheme', ns)
        if not font_scheme:
            return f"Unable to resolve theme code: {theme_code} (no font scheme found)"
            
        # Find major and minor fonts
        major_font = font_scheme.find('.//a:majorFont', ns)
        minor_font = font_scheme.find('.//a:minorFont', ns)
        
        if not major_font or not minor_font:
            return f"Unable to resolve theme code: {theme_code} (incomplete font scheme)"
            
        # Handle theme codes
        # Major/Minor Latin fonts
        if theme_code == "+mj-lt":  # Major font, latin
            latin = major_font.find('.//a:latin', ns)
            return latin.get('typeface') if latin is not None else "Unknown"
        elif theme_code == "+mn-lt":  # Minor font, latin
            latin = minor_font.find('.//a:latin', ns)
            return latin.get('typeface') if latin is not None else "Unknown"
            
        # Major/Minor East Asian fonts
        elif theme_code == "+mj-ea":  # Major font, east asian
            ea = major_font.find('.//a:ea', ns)
            return ea.get('typeface') if ea is not None else "Unknown"
        elif theme_code == "+mn-ea":  # Minor font, east asian
            ea = minor_font.find('.//a:ea', ns)
            return ea.get('typeface') if ea is not None else "Unknown"
            
        # Major/Minor Complex Script fonts
        elif theme_code == "+mj-cs":  # Major font, complex script
            cs = major_font.find('.//a:cs', ns)
            return cs.get('typeface') if cs is not None else "Unknown"
        elif theme_code == "+mn-cs":  # Minor font, complex script
            cs = minor_font.find('.//a:cs', ns)
            return cs.get('typeface') if cs is not None else "Unknown"
            
        # Symbol fonts
        elif theme_code == "+mj-sym":  # Major font, symbol
            sym = major_font.find('.//a:sym', ns)
            return sym.get('typeface') if sym is not None else "Symbol"
        elif theme_code == "+mn-sym":  # Minor font, symbol
            sym = minor_font.find('.//a:sym', ns)
            return sym.get('typeface') if sym is not None else "Symbol"
            
        # Handle other possible theme codes
        elif theme_code.startswith("+mj-"):  # Other major font variants
            script = theme_code[4:]
            return f"Major font for script '{script}' (unresolved)"
        elif theme_code.startswith("+mn-"):  # Other minor font variants
            script = theme_code[4:]
            return f"Minor font for script '{script}' (unresolved)"
        else:
            return f"Unknown theme code: {theme_code}"
    except Exception as e:
        return f"Error resolving theme code '{theme_code}': {str(e)}"

def extract_theme_fonts(presentation: Any) -> Dict[str, Any]:
    """Extract theme font information from the presentation."""
    theme_fonts = {}
    
    try:
        # Get the default theme from the first slide master
        if presentation.slide_masters and len(presentation.slide_masters) > 0:
            master = presentation.slide_masters[0]
            
            # Access the theme through the part relationships
            try:
                # Try to get the theme part through the master's part relationships
                from pptx.opc.constants import RELATIONSHIP_TYPE as RT
                
                # Get the master part - part is on the SlideMaster object, not on _element
                master_part = master.part
                
                # Find theme relationships
                theme_rels = [rel for rel in master_part.rels.values() 
                             if rel.reltype == RT.THEME]
                
                if theme_rels:
                    # Get the first theme part using the relationship ID
                    theme_rel = theme_rels[0]
                    theme_part = master_part.related_part(theme_rel.rId)
                    
                    # Parse the theme XML
                    from xml.etree import ElementTree
                    theme_element = ElementTree.fromstring(theme_part.blob)
                    
                    # Extract font scheme
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    font_scheme_elem = theme_element.find('.//a:fontScheme', ns)
                    
                    if font_scheme_elem is not None:
                        # Get font scheme name
                        scheme_name = font_scheme_elem.get('name', 'Unknown')
                        
                        # Get major font element
                        major_font_elem = font_scheme_elem.find('.//a:majorFont', ns)
                        major_fonts = {}
                        
                        if major_font_elem is not None:
                            latin = major_font_elem.find('.//a:latin', ns)
                            ea = major_font_elem.find('.//a:ea', ns)
                            cs = major_font_elem.find('.//a:cs', ns)
                            
                            major_fonts = {
                                "latin": latin.get('typeface') if latin is not None else None,
                                "east_asian": ea.get('typeface') if ea is not None else None,
                                "complex_script": cs.get('typeface') if cs is not None else None
                            }
                            
                            # Extract detailed font information
                            major_fonts_details = extract_font_details(major_font_elem, ns)
                        
                        # Get minor font element
                        minor_font_elem = font_scheme_elem.find('.//a:minorFont', ns)
                        minor_fonts = {}
                        
                        if minor_font_elem is not None:
                            latin = minor_font_elem.find('.//a:latin', ns)
                            ea = minor_font_elem.find('.//a:ea', ns)
                            cs = minor_font_elem.find('.//a:cs', ns)
                            
                            minor_fonts = {
                                "latin": latin.get('typeface') if latin is not None else None,
                                "east_asian": ea.get('typeface') if ea is not None else None,
                                "complex_script": cs.get('typeface') if cs is not None else None
                            }
                            
                            # Extract detailed font information
                            minor_fonts_details = extract_font_details(minor_font_elem, ns)
                        
                        theme_fonts = {
                            "scheme_name": scheme_name,
                            "major_fonts": major_fonts,
                            "minor_fonts": minor_fonts,
                            "major_fonts_details": major_fonts_details if 'major_fonts_details' in locals() else {},
                            "minor_fonts_details": minor_fonts_details if 'minor_fonts_details' in locals() else {}
                        }
                        
                        # Store the XML for reference
                        theme_fonts["xml"] = ElementTree.tostring(font_scheme_elem).decode()
                
                # Extract default text styles from the slide master
                try:
                    if hasattr(master, '_element'):
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                              'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
                        
                        # Get text styles from the slide master
                        text_styles = master._element.find('.//p:txStyles', ns)
                        if text_styles is not None:
                            theme_fonts["master_text_styles"] = {}
                            
                            # Title style
                            title_style = text_styles.find('.//p:titleStyle', ns)
                            if title_style is not None:
                                theme_fonts["master_text_styles"]["title_style"] = extract_text_style_fonts(title_style, ns, master)
                            
                            # Body style
                            body_style = text_styles.find('.//p:bodyStyle', ns)
                            if body_style is not None:
                                theme_fonts["master_text_styles"]["body_style"] = extract_text_style_fonts(body_style, ns, master)
                            
                            # Other style
                            other_style = text_styles.find('.//p:otherStyle', ns)
                            if other_style is not None:
                                theme_fonts["master_text_styles"]["other_style"] = extract_text_style_fonts(other_style, ns, master)
                except Exception as e:
                    theme_fonts["master_text_styles_error"] = str(e)
            except Exception as e:
                theme_fonts["theme_access_error"] = str(e)
    except Exception as e:
        theme_fonts["error"] = str(e)
    
    return theme_fonts

def extract_font_details(font_elem: Any, ns: Dict[str, str]) -> Dict[str, Any]:
    """Extract detailed font information from a font element."""
    font_details = {}
    
    try:
        # Extract latin font details
        latin = font_elem.find('.//a:latin', ns)
        if latin is not None:
            font_details["latin"] = {
                "typeface": latin.get('typeface'),
                "panose": latin.get('panose'),
                "pitchFamily": latin.get('pitchFamily'),
                "charset": latin.get('charset')
            }
        
        # Extract east asian font details
        ea = font_elem.find('.//a:ea', ns)
        if ea is not None:
            font_details["east_asian"] = {
                "typeface": ea.get('typeface'),
                "panose": ea.get('panose'),
                "pitchFamily": ea.get('pitchFamily'),
                "charset": ea.get('charset')
            }
        
        # Extract complex script font details
        cs = font_elem.find('.//a:cs', ns)
        if cs is not None:
            font_details["complex_script"] = {
                "typeface": cs.get('typeface'),
                "panose": cs.get('panose'),
                "pitchFamily": cs.get('pitchFamily'),
                "charset": cs.get('charset')
            }
        
        # Extract font for specific scripts
        for script_tag in ['a:font', 'a:cs', 'a:ea', 'a:sym']:
            script_fonts = font_elem.findall(f'.//{script_tag}', ns)
            if script_fonts:
                if "script_fonts" not in font_details:
                    font_details["script_fonts"] = []
                
                for font in script_fonts:
                    script = font.get('script')
                    typeface = font.get('typeface')
                    if script and typeface:
                        font_details["script_fonts"].append({
                            "script": script,
                            "typeface": typeface
                        })
    except Exception as e:
        font_details["error"] = str(e)
    
    return font_details

def shape_to_dict(shape: Any) -> Dict[str, Any]:
    """Convert a shape object to a dictionary of its properties."""
    shape_dict = {
        "name": shape.name,
        "shape_type": str(shape.shape_type),
        "width": shape.width,
        "height": shape.height,
        "left": shape.left,
        "top": shape.top,
    }

    # Handle text if present
    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        shape_dict["text_frame"] = {
            "text": shape.text,
            "default_font": None,
            "theme_font": None,
            "resolved_font": None,
            "paragraphs": []
        }

        # Try to get shape-level font defaults from XML
        if hasattr(shape, '_element'):
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            
            # Extract default text style from shape properties
            try:
                shape_style = shape._element.find('.//a:bodyPr', ns)
                if shape_style is not None:
                    shape_dict["text_frame"]["body_properties"] = {
                        "anchor": shape_style.get('anchor'),
                        "wrap_text": shape_style.get('wrap'),
                        "vertical": shape_style.get('vert'),
                        "rotation": shape_style.get('rot')
                    }
            except Exception as e:
                shape_dict["text_frame"]["body_properties_error"] = str(e)
                
            # Look for default paragraph properties
            try:
                default_style = shape._element.find('.//a:lstStyle/a:defPPr', ns)
                if default_style is not None:
                    latin_font = default_style.find('.//a:latin', ns)
                    ea_font = default_style.find('.//a:ea', ns)
                    cs_font = default_style.find('.//a:cs', ns)
                    
                    shape_dict["text_frame"]["default_paragraph_style"] = {
                        "latin_font": {
                            "typeface": latin_font.get('typeface') if latin_font is not None else None,
                            "resolved": resolve_theme_font(shape, latin_font.get('typeface')) if latin_font is not None else None
                        },
                        "east_asian_font": {
                            "typeface": ea_font.get('typeface') if ea_font is not None else None,
                            "resolved": resolve_theme_font(shape, ea_font.get('typeface')) if ea_font is not None else None
                        },
                        "complex_script_font": {
                            "typeface": cs_font.get('typeface') if cs_font is not None else None,
                            "resolved": resolve_theme_font(shape, cs_font.get('typeface')) if cs_font is not None else None
                        }
                    }
            except Exception as e:
                shape_dict["text_frame"]["default_style_error"] = str(e)
                
            # Look for default text run properties
            try:
                default_run_style = shape._element.find('.//a:lstStyle/a:defPPr/a:defRPr', ns)
                if default_run_style is not None:
                    shape_dict["text_frame"]["default_run_style"] = {
                        "size": int(default_run_style.get('sz')) / 100 if default_run_style.get('sz') else None,
                        "bold": default_run_style.get('b') == '1',
                        "italic": default_run_style.get('i') == '1',
                        "underline": default_run_style.get('u') != 'none' if default_run_style.get('u') else False,
                        "strike": default_run_style.get('strike') if default_run_style.get('strike') else None,
                        "baseline": default_run_style.get('baseline') if default_run_style.get('baseline') else None
                    }
                    
                    # Get font information
                    latin_font = default_run_style.find('.//a:latin', ns)
                    ea_font = default_run_style.find('.//a:ea', ns)
                    cs_font = default_run_style.find('.//a:cs', ns)
                    
                    if any([latin_font, ea_font, cs_font]):
                        shape_dict["text_frame"]["default_run_style"]["fonts"] = {
                            "latin": {
                                "typeface": latin_font.get('typeface') if latin_font is not None else None,
                                "resolved": resolve_theme_font(shape, latin_font.get('typeface')) if latin_font is not None else None
                            },
                            "east_asian": {
                                "typeface": ea_font.get('typeface') if ea_font is not None else None,
                                "resolved": resolve_theme_font(shape, ea_font.get('typeface')) if ea_font is not None else None
                            },
                            "complex_script": {
                                "typeface": cs_font.get('typeface') if cs_font is not None else None,
                                "resolved": resolve_theme_font(shape, cs_font.get('typeface')) if cs_font is not None else None
                            }
                        }
            except Exception as e:
                shape_dict["text_frame"]["default_run_style_error"] = str(e)
                
            # Get direct shape-level font properties
            shape_props = shape._element.find('.//a:pPr', ns)
            if shape_props is not None:
                latin_font = shape_props.find('.//a:latin', ns)
                if latin_font is not None:
                    theme_font = latin_font.get('typeface')
                    shape_dict["text_frame"]["theme_font"] = theme_font
                    shape_dict["text_frame"]["resolved_font"] = resolve_theme_font(shape, theme_font)

        # Try to get text frame level defaults
        try:
            if hasattr(shape.text_frame, 'properties'):
                shape_dict["text_frame"]["properties"] = {
                    "margin_left": shape.text_frame.margin_left,
                    "margin_right": shape.text_frame.margin_right,
                    "margin_top": shape.text_frame.margin_top,
                    "margin_bottom": shape.text_frame.margin_bottom,
                    "vertical_anchor": str(shape.text_frame.vertical_anchor),
                    "word_wrap": shape.text_frame.word_wrap,
                    "auto_size": str(shape.text_frame.auto_size) if hasattr(shape.text_frame, 'auto_size') else None
                }
        except Exception as e:
            shape_dict["text_frame"]["properties_error"] = str(e)

        # Process paragraphs
        for p in shape.text_frame.paragraphs:
            para_dict = {
                "text": p.text,
                "level": p.level,
                "alignment": str(p.alignment) if hasattr(p, 'alignment') else None,
                "runs": []
            }

            # Try to get paragraph level font defaults
            try:
                if hasattr(p, 'font'):
                    para_dict["default_font"] = {
                        "name": p.font.name,
                        "size": p.font.size.pt if p.font.size else None,
                        "bold": p.font.bold,
                        "italic": p.font.italic,
                        "underline": p.font.underline
                    }
                # Try to get font info from XML
                if hasattr(p, '_element'):
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    para_props = p._element.find('.//a:pPr', ns)
                    if para_props is not None:
                        latin_font = para_props.find('.//a:latin', ns)
                        ea_font = para_props.find('.//a:ea', ns)
                        cs_font = para_props.find('.//a:cs', ns)
                        
                        if any([latin_font, ea_font, cs_font]):
                            para_dict["theme_fonts"] = {
                                "latin": {
                                    "typeface": latin_font.get('typeface') if latin_font is not None else None,
                                    "resolved": resolve_theme_font(shape, latin_font.get('typeface')) if latin_font is not None else None
                                },
                                "east_asian": {
                                    "typeface": ea_font.get('typeface') if ea_font is not None else None,
                                    "resolved": resolve_theme_font(shape, ea_font.get('typeface')) if ea_font is not None else None
                                },
                                "complex_script": {
                                    "typeface": cs_font.get('typeface') if cs_font is not None else None,
                                    "resolved": resolve_theme_font(shape, cs_font.get('typeface')) if cs_font is not None else None
                                }
                            }
                        
                            # For backward compatibility
                            if latin_font is not None:
                                theme_font = latin_font.get('typeface')
                                para_dict["theme_font"] = theme_font
                                para_dict["resolved_font"] = resolve_theme_font(shape, theme_font)
            except Exception as e:
                para_dict["font_error"] = str(e)

            for run in p.runs:
                run_dict = {
                    "text": run.text,
                    "font": None,
                    "theme_font": None,
                    "resolved_font": None
                }
                try:
                    if hasattr(run, 'font'):
                        run_dict["font"] = {
                            "name": run.font.name,
                            "size": run.font.size.pt if run.font.size else None,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "underline": run.font.underline,
                        }
                    # Try to get run-level font info from XML
                    if hasattr(run, '_element'):
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        run_props = run._element.find('.//a:rPr', ns)
                        if run_props is not None:
                            latin_font = run_props.find('.//a:latin', ns)
                            ea_font = run_props.find('.//a:ea', ns)
                            cs_font = run_props.find('.//a:cs', ns)
                            
                            if any([latin_font, ea_font, cs_font]):
                                run_dict["theme_fonts"] = {
                                    "latin": {
                                        "typeface": latin_font.get('typeface') if latin_font is not None else None,
                                        "resolved": resolve_theme_font(shape, latin_font.get('typeface')) if latin_font is not None else None
                                    },
                                    "east_asian": {
                                        "typeface": ea_font.get('typeface') if ea_font is not None else None,
                                        "resolved": resolve_theme_font(shape, ea_font.get('typeface')) if ea_font is not None else None
                                    },
                                    "complex_script": {
                                        "typeface": cs_font.get('typeface') if cs_font is not None else None,
                                        "resolved": resolve_theme_font(shape, cs_font.get('typeface')) if cs_font is not None else None
                                    }
                                }
                            
                            # For backward compatibility
                            if latin_font is not None:
                                theme_font = latin_font.get('typeface')
                                run_dict["theme_font"] = theme_font
                                run_dict["resolved_font"] = resolve_theme_font(shape, theme_font)
                except Exception as e:
                    run_dict["font_error"] = str(e)
                para_dict["runs"].append(run_dict)

            shape_dict["text_frame"]["paragraphs"].append(para_dict)

    # Handle table if present
    if hasattr(shape, 'has_table') and shape.has_table:
        shape_dict["table"] = {
            "rows": len(shape.table.rows),
            "columns": len(shape.table.columns),
            "cells": [
                [
                    {
                        "text": cell.text,
                        "location": f"row {row_idx}, col {col_idx}"
                    }
                    for col_idx, cell in enumerate(row.cells)
                ]
                for row_idx, row in enumerate(shape.table.rows)
            ]
        }

    # Handle image if present
    if hasattr(shape, 'image'):
        try:
            shape_dict["image"] = {
                "filename": shape.image.filename,
                "content_type": shape.image.content_type,
                "size": shape.image.size,
            }
        except AttributeError:
            shape_dict["image"] = None

    # Add placeholder information if available
    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
        try:
            shape_dict["placeholder"] = {
                "type": str(shape.placeholder_format.type),
                "idx": shape.placeholder_format.idx
            }
        except AttributeError:
            pass

    return shape_dict

def slide_to_dict(slide: Any, slide_index: int) -> Dict[str, Any]:
    """Convert a slide object to a dictionary of its properties."""
    slide_dict = {
        "slide_number": slide_index + 1,
        "shapes": [],
        "slide_id": slide.slide_id,
    }

    # Check if slide is hidden
    if hasattr(slide, '_element'):
        slide_dict["hidden"] = slide._element.get('show') == '0'

    # Get slide layout info
    if hasattr(slide, 'slide_layout'):
        layout_dict = {
            "name": slide.slide_layout.name,
        }
        
        # Try to extract theme font information from the slide layout
        try:
            if hasattr(slide.slide_layout, 'slide_master'):
                master = slide.slide_layout.slide_master
                
                # Access the theme through the part relationships
                try:
                    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
                    
                    # Get the master part - part is on the SlideMaster object, not on _element
                    master_part = master.part
                    
                    # Find theme relationships
                    theme_rels = [rel for rel in master_part.rels.values() 
                                 if rel.reltype == RT.THEME]
                    
                    if theme_rels:
                        # Get the first theme part using the relationship ID
                        theme_rel = theme_rels[0]
                        theme_part = master_part.related_part(theme_rel.rId)
                        
                        # Parse the theme XML
                        from xml.etree import ElementTree
                        theme_element = ElementTree.fromstring(theme_part.blob)
                        
                        # Extract font scheme
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        font_scheme_elem = theme_element.find('.//a:fontScheme', ns)
                        
                        if font_scheme_elem is not None:
                            # Get font scheme name
                            scheme_name = font_scheme_elem.get('name', 'Unknown')
                            
                            # Get major font element
                            major_font_elem = font_scheme_elem.find('.//a:majorFont', ns)
                            major_fonts = {}
                            
                            if major_font_elem is not None:
                                latin = major_font_elem.find('.//a:latin', ns)
                                ea = major_font_elem.find('.//a:ea', ns)
                                cs = major_font_elem.find('.//a:cs', ns)
                                
                                major_fonts = {
                                    "latin": latin.get('typeface') if latin is not None else None,
                                    "east_asian": ea.get('typeface') if ea is not None else None,
                                    "complex_script": cs.get('typeface') if cs is not None else None
                                }
                            
                            # Get minor font element
                            minor_font_elem = font_scheme_elem.find('.//a:minorFont', ns)
                            minor_fonts = {}
                            
                            if minor_font_elem is not None:
                                latin = minor_font_elem.find('.//a:latin', ns)
                                ea = minor_font_elem.find('.//a:ea', ns)
                                cs = minor_font_elem.find('.//a:cs', ns)
                                
                                minor_fonts = {
                                    "latin": latin.get('typeface') if latin is not None else None,
                                    "east_asian": ea.get('typeface') if ea is not None else None,
                                    "complex_script": cs.get('typeface') if cs is not None else None
                                }
                            
                            layout_dict["theme_fonts"] = {
                                "scheme_name": scheme_name,
                                "major_fonts": major_fonts,
                                "minor_fonts": minor_fonts
                            }
                    
                    # Try to extract default text styles from the slide master
                    try:
                        if hasattr(master, '_element'):
                            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                                  'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
                            
                            # Get text styles from the slide master
                            text_styles = master._element.find('.//p:txStyles', ns)
                            if text_styles is not None:
                                layout_dict["master_text_styles"] = {}
                                
                                # Title style
                                title_style = text_styles.find('.//p:titleStyle', ns)
                                if title_style is not None:
                                    layout_dict["master_text_styles"]["title_style"] = extract_text_style_fonts(title_style, ns, slide)
                                
                                # Body style
                                body_style = text_styles.find('.//p:bodyStyle', ns)
                                if body_style is not None:
                                    layout_dict["master_text_styles"]["body_style"] = extract_text_style_fonts(body_style, ns, slide)
                                
                                # Other style
                                other_style = text_styles.find('.//p:otherStyle', ns)
                                if other_style is not None:
                                    layout_dict["master_text_styles"]["other_style"] = extract_text_style_fonts(other_style, ns, slide)
                    except Exception as e:
                        layout_dict["master_text_styles_error"] = str(e)
                except Exception as e:
                    layout_dict["theme_access_error"] = str(e)
        except Exception as e:
            layout_dict["theme_fonts_error"] = str(e)
            
        slide_dict["layout"] = layout_dict

    # Get background info if available
    try:
        background = slide._element.find('.//p:bg',
            {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if background is not None:
            slide_dict["background"] = ElementTree.tostring(background).decode()
    except Exception:
        slide_dict["background"] = None

    # Check for transitions
    try:
        transition = slide._element.find('.//p:transition',
            {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if transition is not None:
            slide_dict["has_transition"] = True
            slide_dict["transition_xml"] = ElementTree.tostring(transition).decode()
    except Exception:
        slide_dict["has_transition"] = False

    # Check for animations
    try:
        timing = slide._element.find('.//p:timing',
            {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if timing is not None:
            slide_dict["has_animations"] = True
            slide_dict["timing_xml"] = ElementTree.tostring(timing).decode()
    except Exception:
        slide_dict["has_animations"] = False

    # Convert each shape
    for shape in slide.shapes:
        try:
            shape_dict = shape_to_dict(shape)
            slide_dict["shapes"].append(shape_dict)
        except Exception as e:
            slide_dict["shapes"].append({
                "error": f"Failed to convert shape: {str(e)}",
                "shape_name": getattr(shape, 'name', 'unknown')
            })

    return slide_dict

def extract_text_style_fonts(style_element: Any, ns: Dict[str, str], slide: Any) -> Dict[str, Any]:
    """Extract font information from a text style element."""
    result = {}
    
    try:
        # Get default paragraph properties
        def_p_pr = style_element.find('.//a:defPPr', ns)
        if def_p_pr is not None:
            result["default_paragraph"] = {}
            
            # Get default run properties
            def_r_pr = def_p_pr.find('.//a:defRPr', ns)
            if def_r_pr is not None:
                result["default_paragraph"]["default_run"] = {
                    "size": int(def_r_pr.get('sz')) / 100 if def_r_pr.get('sz') else None,
                    "bold": def_r_pr.get('b') == '1',
                    "italic": def_r_pr.get('i') == '1',
                    "underline": def_r_pr.get('u') != 'none' if def_r_pr.get('u') else False,
                }
                
                # Get font information
                latin_font = def_r_pr.find('.//a:latin', ns)
                ea_font = def_r_pr.find('.//a:ea', ns)
                cs_font = def_r_pr.find('.//a:cs', ns)
                
                if any([latin_font, ea_font, cs_font]):
                    result["default_paragraph"]["default_run"]["fonts"] = {
                        "latin": {
                            "typeface": latin_font.get('typeface') if latin_font is not None else None,
                            "resolved": resolve_theme_font(slide, latin_font.get('typeface')) if latin_font is not None else None
                        },
                        "east_asian": {
                            "typeface": ea_font.get('typeface') if ea_font is not None else None,
                            "resolved": resolve_theme_font(slide, ea_font.get('typeface')) if ea_font is not None else None
                        },
                        "complex_script": {
                            "typeface": cs_font.get('typeface') if cs_font is not None else None,
                            "resolved": resolve_theme_font(slide, cs_font.get('typeface')) if cs_font is not None else None
                        }
                    }
        
        # Get level paragraph properties (for different outline levels)
        level_p_prs = []
        for i in range(1, 10):  # Check levels 1-9
            level_p_pr = style_element.find(f'.//a:lvl{i}pPr', ns)
            if level_p_pr is not None:
                level_p_prs.append((i, level_p_pr))
        
        if level_p_prs:
            result["levels"] = {}
            
            for level_idx, lvl_p_pr in level_p_prs:
                result["levels"][f"level_{level_idx}"] = {}
                
                # Get run properties for this level
                r_pr = lvl_p_pr.find('.//a:defRPr', ns)
                if r_pr is not None:
                    result["levels"][f"level_{level_idx}"]["run_properties"] = {
                        "size": int(r_pr.get('sz')) / 100 if r_pr.get('sz') else None,
                        "bold": r_pr.get('b') == '1',
                        "italic": r_pr.get('i') == '1',
                        "underline": r_pr.get('u') != 'none' if r_pr.get('u') else False,
                    }
                    
                    # Get font information
                    latin_font = r_pr.find('.//a:latin', ns)
                    ea_font = r_pr.find('.//a:ea', ns)
                    cs_font = r_pr.find('.//a:cs', ns)
                    
                    if any([latin_font, ea_font, cs_font]):
                        result["levels"][f"level_{level_idx}"]["run_properties"]["fonts"] = {
                            "latin": {
                                "typeface": latin_font.get('typeface') if latin_font is not None else None,
                                "resolved": resolve_theme_font(slide, latin_font.get('typeface')) if latin_font is not None else None
                            },
                            "east_asian": {
                                "typeface": ea_font.get('typeface') if ea_font is not None else None,
                                "resolved": resolve_theme_font(slide, ea_font.get('typeface')) if ea_font is not None else None
                            },
                            "complex_script": {
                                "typeface": cs_font.get('typeface') if cs_font is not None else None,
                                "resolved": resolve_theme_font(slide, cs_font.get('typeface')) if cs_font is not None else None
                            }
                        }
    except Exception as e:
        result["error"] = str(e)
    
    return result

def presentation_to_dict(pptx_path: Path) -> Dict[str, Any]:
    """Convert entire presentation to a dictionary."""
    prs = Presentation(pptx_path)

    # Basic presentation properties
    pres_dict = {
        "metadata": {
            "slides_count": len(prs.slides),
            "core_properties": {
                "author": prs.core_properties.author,
                "created": str(prs.core_properties.created) if prs.core_properties.created else None,
                "modified": str(prs.core_properties.modified) if prs.core_properties.modified else None,
                "title": prs.core_properties.title,
                "subject": prs.core_properties.subject,
                "keywords": prs.core_properties.keywords,
                "comments": prs.core_properties.comments,
                "category": prs.core_properties.category,
            },
            "theme_fonts": extract_theme_fonts(prs)
        },
        "slides": []
    }

    # Convert each slide
    for idx, slide in enumerate(prs.slides):
        try:
            slide_dict = slide_to_dict(slide, idx)
            pres_dict["slides"].append(slide_dict)
        except Exception as e:
            pres_dict["slides"].append({
                "error": f"Failed to convert slide {idx + 1}: {str(e)}",
                "slide_number": idx + 1
            })

    return pres_dict

def main():
    if len(sys.argv) != 2:
        print("Usage: python script.py <path_to_pptx>")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"Error: File '{pptx_path}' does not exist")
        sys.exit(1)

    if pptx_path.suffix.lower() != '.pptx':
        print("Error: File must be a .pptx file")
        sys.exit(1)

    try:
        pres_dict = presentation_to_dict(pptx_path)
        print(json.dumps(pres_dict, indent=2, default=str))
    except Exception as e:
        print(f"Error processing presentation: {str(e)}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
